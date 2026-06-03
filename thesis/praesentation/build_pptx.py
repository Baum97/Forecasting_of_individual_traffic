# -*- coding: utf-8 -*-
"""
Generiert die Verteidigungs-Praesentation (verteidigung.pptx) aus dem
inhaltlichen Stand der Masterarbeit.

Quelle der Inhalte:
  - thesis/chapters/01einleitung.tex  (Motivation, Problem, Ziele, Fragen)
  - thesis/chapters/02grundlagen.tex  (Forschungsstraenge, Fahrertypen)
  - thesis/chapters/todo.tex          (Aenderungslog = implementiertes System)
  - predictions/cross_validation_matrix.csv (echte Ergebnisse)

Aufruf (aus Repo-Root):
  python thesis/praesentation/build_pptx.py
"""

import csv
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------
# Pfade
# --------------------------------------------------------------------------
HERE = os.path.dirname(os.path.abspath(__file__))
THESIS = os.path.dirname(HERE)
ROOT = os.path.dirname(THESIS)
IMAGES = os.path.join(THESIS, "images")
PRED = os.path.join(ROOT, "predictions")
ASSETS = os.path.join(HERE, "assets")
os.makedirs(ASSETS, exist_ok=True)
OUT = os.path.join(HERE, "verteidigung.pptx")

LOGO_HS = os.path.join(IMAGES, "HochschuleEsslingen_Logo_4c_DE.png")

# --------------------------------------------------------------------------
# Farbpalette (an HS Esslingen / sachlich-technisch angelehnt)
# --------------------------------------------------------------------------
NAVY = RGBColor(0x10, 0x2A, 0x43)      # dunkles Navy (Titel/Flaechen)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)      # Akzent-Blau
TEAL = RGBColor(0x17, 0xA2, 0x9A)      # Akzent-Teal
ORANGE = RGBColor(0xE8, 0x7D, 0x1E)    # Highlight
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)     # heller Hintergrund
MIDGRAY = RGBColor(0x5A, 0x66, 0x72)   # Fliesstext-grau
DARK = RGBColor(0x1B, 0x21, 0x28)      # fast schwarz
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_LIGHT = "Calibri Light"

# 16:9
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# --------------------------------------------------------------------------
# Hilfsfunktionen
# --------------------------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _set_fill(sp, color)
    if line is not None:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_para(tf, text, size, color, bold=False, italic=False, font=FONT,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0, bullet=None,
             level=0, line_spacing=None):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    if line_spacing:
        p.line_spacing = line_spacing
    if bullet is not None:
        r0 = p.add_run()
        set_run(r0, bullet + "  ", size, color, bold=True, font=font)
    r = p.add_run()
    set_run(r, text, size, color, bold=bold, italic=italic, font=font)
    return p


# --------------------------------------------------------------------------
# Folien-Bausteine
# --------------------------------------------------------------------------
PAGE = {"n": 0}


def content_slide(title, kicker=None):
    """Standard-Inhaltsfolie: Kopfband + Titel. Gibt slide zurueck."""
    PAGE["n"] += 1
    slide = prs.slides.add_slide(BLANK)
    # heller Hintergrund
    rect(slide, 0, 0, EMU_W, EMU_H, WHITE)
    # linker Akzentbalken
    rect(slide, 0, 0, Inches(0.18), EMU_H, BLUE)
    # Titelzone
    tb, tf = textbox(slide, Inches(0.6), Inches(0.42), Inches(11.6), Inches(1.1))
    if kicker:
        add_para(tf, kicker.upper(), 12, TEAL, bold=True, space_after=2)
    add_para(tf, title, 28, NAVY, bold=True, font=FONT_LIGHT, space_after=0)
    # Trennlinie
    rect(slide, Inches(0.62), Inches(1.5), Inches(11.9), Pt(2.2), BLUE)
    # Fusszeile
    fb, ff = textbox(slide, Inches(0.6), Inches(7.06), Inches(11.9), Inches(0.32),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = ff.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, "Prognose individuellen Mobilitätsverhaltens  ·  Achim Baumgärtner", 9, MIDGRAY)
    # Seitenzahl
    nb, nf = textbox(slide, Inches(12.4), Inches(7.06), Inches(0.7), Inches(0.32),
                     anchor=MSO_ANCHOR.MIDDLE)
    pp = nf.paragraphs[0]
    pp.alignment = PP_ALIGN.RIGHT
    rr = pp.add_run()
    set_run(rr, str(PAGE["n"]), 9, MIDGRAY)
    return slide


def body_box(slide, x=Inches(0.62), y=Inches(1.75), w=Inches(11.9), h=Inches(5.0)):
    tb, tf = textbox(slide, x, y, w, h)
    return tf


def section_divider(number, title, subtitle=None):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, EMU_W, EMU_H, NAVY)
    rect(slide, 0, Inches(3.05), EMU_W, Inches(1.4), BLUE)
    tb, tf = textbox(slide, Inches(0.9), Inches(3.15), Inches(11.5), Inches(1.2),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, f"{number}", 20, RGBColor(0xBF, 0xD8, 0xEE), bold=True, space_after=2)
    add_para(tf, title, 34, WHITE, bold=True, font=FONT_LIGHT, space_after=0)
    if subtitle:
        sb, sf = textbox(slide, Inches(0.92), Inches(4.6), Inches(11.3), Inches(0.8))
        add_para(sf, subtitle, 14, RGBColor(0xBF, 0xD8, 0xEE))
    return slide


# --------------------------------------------------------------------------
# Cross-Validation-Heatmap (F1) aus CSV erzeugen
# --------------------------------------------------------------------------
def build_cv_heatmap():
    rows = {}
    with open(os.path.join(PRED, "cross_validation_matrix.csv"), encoding="utf-8") as f:
        for r in csv.DictReader(f):
            rows[(r["model_trained_on"], r["evaluated_on"])] = r

    order = ["routine", "real_world_ev", "ved", "emobpy"]
    labels = {"routine": "routine", "real_world_ev": "real_world_ev",
              "ved": "ved", "emobpy": "emobpy"}
    n = len(order)
    F1 = np.zeros((n, n))
    AUC = np.zeros((n, n))
    for i, tr in enumerate(order):
        for j, ev in enumerate(order):
            F1[i, j] = float(rows[(tr, ev)]["f1"])
            AUC[i, j] = float(rows[(tr, ev)]["roc_auc"])

    def _plot(M, fname, title, cmap, fmt):
        fig, ax = plt.subplots(figsize=(6.2, 5.2), dpi=200)
        im = ax.imshow(M, cmap=cmap, vmin=0, vmax=max(0.9, M.max()))
        ax.set_xticks(range(n)); ax.set_yticks(range(n))
        ax.set_xticklabels([labels[o] for o in order], rotation=30, ha="right", fontsize=10)
        ax.set_yticklabels([labels[o] for o in order], fontsize=10)
        ax.set_xlabel("evaluiert auf", fontsize=11, fontweight="bold")
        ax.set_ylabel("trainiert auf", fontsize=11, fontweight="bold")
        ax.set_title(title, fontsize=13, fontweight="bold", pad=12, color="#102A43")
        for i in range(n):
            for j in range(n):
                v = M[i, j]
                ax.text(j, i, fmt.format(v),
                        ha="center", va="center", fontsize=11,
                        color="white" if v > 0.55 * M.max() else "#102A43",
                        fontweight="bold")
        # Diagonale markieren
        for i in range(n):
            ax.add_patch(plt.Rectangle((i - 0.5, i - 0.5), 1, 1, fill=False,
                                       edgecolor="#E87D1E", lw=2.5))
        fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
        fig.tight_layout()
        out = os.path.join(ASSETS, fname)
        fig.savefig(out, bbox_inches="tight")
        plt.close(fig)
        return out

    f1_png = _plot(F1, "cv_f1.png", "F1-Score  (4×4 Cross-Validation)", "viridis", "{:.2f}")
    auc_png = _plot(AUC, "cv_auc.png", "ROC-AUC  (4×4 Cross-Validation)", "cividis", "{:.2f}")
    return f1_png, auc_png, F1, AUC, order


# --------------------------------------------------------------------------
# Pipeline-Diagramm (Feature -> 2 Modellfamilien)
# --------------------------------------------------------------------------
def build_pipeline_diagram():
    fig, ax = plt.subplots(figsize=(10.5, 3.6), dpi=200)
    ax.axis("off")
    ax.set_xlim(0, 100); ax.set_ylim(0, 40)

    def box(x, y, w, h, text, fc, tc="white", fs=10, bold=True):
        ax.add_patch(plt.Rectangle((x, y), w, h, facecolor=fc, edgecolor="none",
                                   zorder=2, joinstyle="round"))
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center",
                color=tc, fontsize=fs, fontweight="bold" if bold else "normal",
                zorder=3, wrap=True)

    def arrow(x1, y1, x2, y2):
        ax.annotate("", xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle="-|>", color="#5A6672", lw=2))

    # Quellen
    box(1, 28, 20, 9, "4 Datenquellen\nemobpy · real_world_ev\nved · routine", "#1F6FB2", fs=9)
    # Adapter
    box(1, 15, 20, 8, "data_adapters.py\neinheitliches Schema\n(id, timestamp, driving)", "#17A29A", fs=9)
    arrow(11, 28, 11, 23)
    # Feature Engineering
    box(28, 18, 22, 13, "Feature-Engineering\nhour · weekday · is_weekend\nLags 1/2/24/168\nRolling 24h / 168h", "#102A43", fs=9)
    arrow(21, 19, 28, 24)
    # Zwei Familien
    box(58, 27, 40, 10, "Binärer Driving-Classifier\nRandomForest · 4×4-CV-Matrix\n(Datenqualität vs. Algorithmus)", "#E87D1E", fs=9)
    box(58, 12, 40, 11, "Forecaster (1–7 Tage)\nRF / LightGBM · Monte-Carlo\nP10/P90-Quantile · Trip-Blöcke", "#1F6FB2", fs=9)
    arrow(50, 26, 58, 31)
    arrow(50, 23, 58, 18)
    fig.tight_layout()
    out = os.path.join(ASSETS, "pipeline.png")
    fig.savefig(out, bbox_inches="tight", transparent=True)
    plt.close(fig)
    return out


# --------------------------------------------------------------------------
# Bullet-Folie (einspaltig)
# --------------------------------------------------------------------------
def bullets_slide(title, kicker, items, sizes=18):
    slide = content_slide(title, kicker)
    tf = body_box(slide)
    for it in items:
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        color = DARK if lvl == 0 else MIDGRAY
        bullet = "▸" if lvl == 0 else "–"
        sz = sizes if lvl == 0 else sizes - 3
        add_para(tf, txt, sz, color, bullet=bullet, level=lvl,
                 space_after=9 if lvl == 0 else 5, line_spacing=1.05)
    return slide


# --------------------------------------------------------------------------
# Zwei-Spalten-Folie
# --------------------------------------------------------------------------
def two_col_slide(title, kicker, left_head, left_items, right_head, right_items):
    slide = content_slide(title, kicker)
    # linke Karte
    rect(slide, Inches(0.62), Inches(1.8), Inches(5.75), Inches(4.9), LIGHT)
    rect(slide, Inches(0.62), Inches(1.8), Inches(5.75), Inches(0.55), BLUE)
    lh, lhf = textbox(slide, Inches(0.85), Inches(1.86), Inches(5.3), Inches(0.45),
                      anchor=MSO_ANCHOR.MIDDLE)
    add_para(lhf, left_head, 15, WHITE, bold=True)
    lt, ltf = textbox(slide, Inches(0.9), Inches(2.55), Inches(5.2), Inches(4.0))
    for it in left_items:
        add_para(ltf, it, 14, DARK, bullet="▸", space_after=8, line_spacing=1.03)
    # rechte Karte
    rect(slide, Inches(6.95), Inches(1.8), Inches(5.75), Inches(4.9), LIGHT)
    rect(slide, Inches(6.95), Inches(1.8), Inches(5.75), Inches(0.55), TEAL)
    rh, rhf = textbox(slide, Inches(7.18), Inches(1.86), Inches(5.3), Inches(0.45),
                      anchor=MSO_ANCHOR.MIDDLE)
    add_para(rhf, right_head, 15, WHITE, bold=True)
    rt, rtf = textbox(slide, Inches(7.23), Inches(2.55), Inches(5.2), Inches(4.0))
    for it in right_items:
        add_para(rtf, it, 14, DARK, bullet="▸", space_after=8, line_spacing=1.03)
    return slide


# ==========================================================================
#  Assets bauen
# ==========================================================================
f1_png, auc_png, F1, AUC, cv_order = build_cv_heatmap()
pipe_png = build_pipeline_diagram()


# ==========================================================================
#  FOLIE 1 — Titel
# ==========================================================================
def title_slide():
    PAGE["n"] = 0
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, EMU_W, EMU_H, NAVY)
    rect(slide, 0, 0, Inches(0.25), EMU_H, ORANGE)
    rect(slide, 0, Inches(4.55), EMU_W, Pt(2.5), BLUE)
    # Logo
    if os.path.exists(LOGO_HS):
        slide.shapes.add_picture(LOGO_HS, Inches(0.7), Inches(0.6), height=Inches(0.7))
    # Label
    tb, tf = textbox(slide, Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.5))
    add_para(tf, "MASTERARBEIT  ·  VERTEIDIGUNG", 14, TEAL, bold=True)
    # Titel
    tt, ttf = textbox(slide, Inches(0.75), Inches(2.55), Inches(11.9), Inches(1.9))
    add_para(ttf, "Prognose individuellen Mobilitätsverhaltens",
             36, WHITE, bold=True, font=FONT_LIGHT, space_after=2, line_spacing=1.0)
    add_para(ttf, "auf Basis eines mit Bewegungsdaten trainierten Machine-Learning-Modells",
             20, RGBColor(0xBF, 0xD8, 0xEE), font=FONT_LIGHT, line_spacing=1.0)
    # Autor / Gutachter
    ab, af = textbox(slide, Inches(0.75), Inches(4.9), Inches(11.8), Inches(1.8))
    add_para(af, "Achim Baumgärtner", 20, WHITE, bold=True, space_after=4)
    add_para(af, "Matrikelnummer 766010   ·   Fakultät Informatik, Hochschule Esslingen",
             13, RGBColor(0xCE, 0xD9, 0xE3), space_after=10)
    add_para(af, "Erstgutachter: Prof. Dr. Sonntag      Zweitgutachter: Prof. Dr. Schober",
             13, RGBColor(0xCE, 0xD9, 0xE3))


title_slide()

# ==========================================================================
#  FOLIE 2 — Agenda
# ==========================================================================
ag = content_slide("Agenda", "Überblick")
left = ["1.  Motivation & Problemstellung", "2.  Zielsetzung & Forschungsfragen",
        "3.  Stand der Forschung", "4.  Datengrundlage"]
right = ["5.  Methodik & System", "6.  Ergebnisse: Datenqualität vs. Algorithmus",
         "7.  Diskussion & Limitationen", "8.  Fazit & Ausblick"]
rect(ag, Inches(0.62), Inches(2.0), Inches(5.75), Inches(4.3), LIGHT)
rect(ag, Inches(6.95), Inches(2.0), Inches(5.75), Inches(4.3), LIGHT)
_, af1 = textbox(ag, Inches(1.0), Inches(2.35), Inches(5.2), Inches(3.7))
for t in left:
    add_para(af1, t, 17, NAVY, bold=True, space_after=18)
_, af2 = textbox(ag, Inches(7.35), Inches(2.35), Inches(5.2), Inches(3.7))
for t in right:
    add_para(af2, t, 17, NAVY, bold=True, space_after=18)

# ==========================================================================
#  Abschnitt 1
# ==========================================================================
section_divider("01", "Motivation & Problemstellung",
                "Warum individuelle Mobilitätsprognose für die Energiewende zählt")

# Motivation
bullets_slide(
    "Vom Verbraucher zum mobilen Speicher", "Motivation",
    [
        "Elektromobilität reduziert Emissionen und fossile Abhängigkeiten – und eröffnet eine neue Rolle in der Energieversorgung.",
        ("Vehicle-to-Grid (V2G): Fahrzeuge nehmen Energie nicht nur auf, sondern geben sie ans Netz zurück.", 1),
        ("Mobile Akkus könnten überschüssige erneuerbare Energie speichern – ein zentraler Baustein der Energiewende.", 1),
        "Der Nutzen entsteht aber nur bei intelligenter Steuerung der Lade-/Entladezyklen.",
        ("Steuerung setzt voraus, dass man weiß, wann ein Fahrzeug fährt und wann es geparkt (= verfügbar) ist.", 1),
        ("Fehlprognosen → fehlerhafte Energieverwaltung, im Extremfall Nichtverfügbarkeit des Fahrzeugs.", 1),
    ])

# Problemstellung / Forschungsluecke
bullets_slide(
    "Forschungslücke: das Individuum", "Problemstellung",
    [
        "Bestehende Arbeiten prognostizieren v. a. Ladezeitpunkt, Leistung, Ladezustand – und aggregiertes Verhalten ganzer Gruppen.",
        "Das individuelle Fahr-/Nutzungsverhalten einzelner Fahrzeuge wird kaum modelliert – gerade dieses ist für die Steuerung nötig.",
        "Individualverkehr ist stark heterogen: ein einzelnes übergreifendes Modell erfasst den Kontext nicht vollständig.",
        "Datenbasis der wenigen einschlägigen Studien ist oft unzureichend oder – bei synthetischen Profilen – nicht zuverlässig repräsentativ.",
        ("Kernproblem: stabile, hochqualitative Datengrundlage – wobei „Qualität\" hier erst über Kriterien definiert werden muss.", 1),
    ])

# ==========================================================================
#  Abschnitt 2 — Ziel & Forschungsfragen
# ==========================================================================
section_divider("02", "Zielsetzung & Forschungsfragen")

bullets_slide(
    "Zielsetzung", "Was die Arbeit leisten soll",
    [
        "Zuverlässige Modelle zur Prognose individuellen Mobilitätsverhaltens entwickeln.",
        ("Konkret: pro Stunde vorhersagen, wann ein Fahrzeug fährt und wann es parkt – über einen Horizont von bis zu 7 Tagen.", 1),
        "Datengrundlage qualitativ bewerten und, wo nötig, aus Datenschutzgründen anonymisieren.",
        "Vorhersagegenauigkeit verschiedener Modelle vergleichen und Einflussfaktoren quantifizieren.",
        "Beitrag: Grundlage für eine intelligente Steuerung der Energieflüsse zwischen Fahrzeug und Netz.",
    ])

# Forschungsfragen als gestaltete Folie
fq = content_slide("Forschungsfragen", "Leitfaden der Arbeit")
cards = [
    ("FF1", BLUE, "Wie zuverlässig lässt sich individuelles Mobilitätsverhalten einzelner Fahrzeuge über bis zu 7 Tage vorhersagen?"),
    ("FF1a", TEAL, "Wie unterscheiden sich ML-Verfahren in ihrer Vorhersagegenauigkeit?"),
    ("FF1a-i", TEAL, "Welchen Einfluss hat die Datenqualität / -quelle – getrennt vom Algorithmus-Effekt?"),
    ("FF1b", TEAL, "Welche Merkmale beeinflussen das Verhalten am stärksten?"),
    ("FF2", ORANGE, "Welche Anforderungen an die Datengrundlage ergeben sich, und wie bewertet man ihre Qualität?"),
    ("FF2a", ORANGE, "Wie wirken Aufbereitung und Anonymisierung auf die Nutzbarkeit?"),
]
ys = [1.8, 2.75, 3.7, 4.65, 5.6, 6.55]
for (tag, col, txt), y in zip(cards, ys):
    rect(fq, Inches(0.62), Inches(y), Inches(1.4), Inches(0.82), col)
    _, tcf = textbox(fq, Inches(0.62), Inches(y), Inches(1.4), Inches(0.82),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tcf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); set_run(r, tag, 14, WHITE, bold=True)
    _, bcf = textbox(fq, Inches(2.2), Inches(y), Inches(10.3), Inches(0.82),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(bcf, txt, 14, DARK, space_after=0, line_spacing=1.0)

# ==========================================================================
#  Abschnitt 3 — Stand der Forschung
# ==========================================================================
section_divider("03", "Stand der Forschung")

# Forschungsstraenge Tabelle
fs = content_slide("Vier Forschungsstränge", "Einordnung")
hdr = ["Forschungsstrang", "Typische Methoden", "Offener Punkt"]
data = [
    ["Mobilitäts-/Nutzungsforecast", "ARIMA, Random Forest, LSTM", "Generalisierung über Fahrzeuge & Quellen"],
    ["Driving-/Parking-Klassifikation", "Random Forest, SVM, Gradient Boosting", "Einfluss der Datenqualität auf die Güte"],
    ["EV-Last- & Ladeprognose", "Gradient Boosting, Deep Learning", "Vergleichbarkeit synthetisch vs. real"],
    ["Unsicherheits-/Quantilsprognose", "Quantilsregression, Quantile-RF", "Kalibrierung individueller P10/P90-Intervalle"],
]
tbl_x, tbl_y = Inches(0.62), Inches(1.9)
col_w = [Inches(3.5), Inches(4.3), Inches(4.1)]
row_h = Inches(0.95)
# Header
cx = tbl_x
for j, h in enumerate(hdr):
    rect(fs, cx, tbl_y, col_w[j], Inches(0.55), NAVY)
    _, hf = textbox(fs, cx + Inches(0.12), tbl_y, col_w[j] - Inches(0.24), Inches(0.55),
                    anchor=MSO_ANCHOR.MIDDLE)
    add_para(hf, h, 13, WHITE, bold=True, space_after=0)
    cx += col_w[j]
# Rows
for i, row in enumerate(data):
    ry = tbl_y + Inches(0.55) + i * row_h
    bg = LIGHT if i % 2 == 0 else WHITE
    cx = tbl_x
    for j, cell in enumerate(row):
        rect(fs, cx, ry, col_w[j], row_h, bg)
        _, cf = textbox(fs, cx + Inches(0.12), ry, col_w[j] - Inches(0.24), row_h,
                        anchor=MSO_ANCHOR.MIDDLE)
        add_para(cf, cell, 12.5, DARK if j > 0 else NAVY,
                 bold=(j == 0), space_after=0, line_spacing=1.0)
        cx += col_w[j]
# Highlight zweite Zeile (eigener Fokus)
hl = fs.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_x, tbl_y + Inches(0.55) + row_h,
                         sum(col_w, Emu(0)), row_h)
hl.fill.background(); hl.line.color.rgb = ORANGE; hl.line.width = Pt(2.5)
hl.shadow.inherit = False

# Fahrertypen
bullets_slide(
    "Heterogenität: fünf Fahrertypen", "ji2023rethinking · >4 Mio. Fahrten, 3743 Fahrer",
    [
        "Homebodies – Wenigfahrer, kurze und seltene Wege, geringe Mobilität.",
        "Gig Drivers – hohes, auftragsgetriebenes Aufkommen ohne festen Arbeitsort.",
        "Movers – breit gestreutes Verhalten, viele wechselnde Ziele, kein klares Muster.",
        "Typical Drivers – Mix aus Arbeit & Freizeit, mittlere Regelmäßigkeit.",
        "Work-focused Commuters – stark regelmäßiges Pendelmuster, hohe Vorhersagbarkeit.",
        ("Konsequenz: ein einziges Modell wird dem Individualverkehr nicht gerecht → individuelle / quellenspezifische Modellierung.", 0),
    ])

# ==========================================================================
#  Abschnitt 4 — Datengrundlage
# ==========================================================================
section_divider("04", "Datengrundlage")

two_col_slide(
    "Vier Datenquellen, ein Schema", "Datengrundlage",
    "Quellen (via data_adapters.py)",
    ["emobpy – synthetisch generierte EV-Profile",
     "real_world_ev – reale EV-Flottendaten",
     "ved – Vehicle Energy Dataset",
     "routine – reale, regelmäßige Nutzungsprofile"],
    "Einheitliches Schema",
    ["vehicle_id · timestamp · driving",
     "Stündliche Auflösung (floor zur vollen Stunde)",
     "Stark unterschiedliche Aktiv-Raten: ~2 % (real_world_ev) bis ~26 % (routine)",
     "Klassen-Imbalance ist die zentrale Herausforderung"])

# ==========================================================================
#  Abschnitt 5 — Methodik & System
# ==========================================================================
section_divider("05", "Methodik & System")

# Pipeline-Diagramm
pp = content_slide("Zwei Modellfamilien, ein Feature-Kern", "Architektur")
pp.shapes.add_picture(pipe_png, Inches(0.7), Inches(1.85), width=Inches(12.0))
_, pcf = textbox(pp, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5))
add_para(pcf, "Bewusste Trennung der Forschungsschienen: die Cross-Validation (Datenqualität, binär) "
              "ist von der Anwendungsschiene (Tages-Forecast, Quantile) entkoppelt.",
         12.5, MIDGRAY, italic=True)

# 4x4 CV Design
bullets_slide(
    "Studiendesign: die 4×4-Matrix", "Datenqualität isolieren",
    [
        "Idee: Den Effekt der Datenquelle vom Effekt des Algorithmus trennen (FF1a-i).",
        "Jedes der 4 Quellen-Modelle wird auf jeder der 4 Quellen evaluiert → 16 Kombinationen.",
        ("Diagonale = trainiert und evaluiert auf derselben Quelle (Best Case).", 1),
        ("Off-Diagonale = Generalisierung über Quellen hinweg (Transfer).", 1),
        "Binärer RandomForest-Classifier „fährt vs. geparkt\" pro Stunde, identische Features für alle Quellen.",
        "Metriken: Accuracy, F1, Precision/Recall, ROC-AUC – F1 wegen starker Imbalance im Fokus.",
    ])

# System / Implementierung
two_col_slide(
    "Implementiertes System", "Backend · Frontend",
    "Backend (FastAPI, Python 3.13)",
    ["Endpunkte: /api/models, /api/forecast, /api/runs, /api/train",
     "Forecaster-Registry: name → Klasse (Duck Typing, keine Basisklasse)",
     "Persistenz: *.joblib pro Modell, Run-Verzeichnisse pro Simulation",
     "Ablations-fähig: RF und LightGBM austauschbar"],
    "Frontend (Angular 17+, Signals)",
    ["Forecast-Tab: Tageskarten + Stunden-Heatmap (Tag × Stunde)",
     "Training-Tab: Läufe ohne Terminal starten, Live-Log",
     "Trip-Blöcke pro Tag (Abfahrt/Rückkehr, P10/P90)",
     "MC- vs. Soft-Mode-Toggle für die Ablation"])

# Hourly Forecasting + Quantile
bullets_slide(
    "Auto-regressives Monte-Carlo-Forecasting", "Wann, nicht nur ob",
    [
        "Stündliches Schwestermodell rollt Stunde für Stunde vorwärts (Lags 1/2/24/168, Rolling 24h/168h).",
        "100 Monte-Carlo-Pfade ziehen je Stunde Bernoulli(p) und reichen das Ergebnis als lag_1 weiter.",
        ("Unsicherheit fällt aus der Sample-Streuung heraus → P10/P50/P90 statt separater Annahme.", 1),
        ("Breite P10/P90-Bänder zeigen explizit, wann der Forecast nicht vertrauenswürdig ist.", 1),
        "Hybrid: Stunden-Grid ist die Wahrheit, Tageskarten werden daraus aggregiert.",
    ])

# Forecast-Kollaps + Fixes (technische Tiefe fuer Q&A)
bullets_slide(
    "Problem & Lösung: Forecast-Kollaps", "Methodischer Befund",
    [
        "Befund: bei niedriger Basisrate (2–10 % aktiv) kollabierte der Forecast auf pUsed = 0.",
        ("Ursache: Klassen-Imbalance + selbstverstärkender autoregressiver Kollaps (frühe Nullen senken p weiter).", 1),
        "Drei Eingriffe:",
        ("class_weight = \"balanced\" – hebt die Minderheitsklasse in den kalibrierten Bereich.", 1),
        ("Adaptive Schwelle p ≥ clamp(1.5 · base_rate, 0.30, 0.60) statt fixem 0.5.", 1),
        ("Soft-Mode – deterministischer Rollout als saubere Ablation gegen Sampling-Rauschen.", 1),
    ])

# ==========================================================================
#  Abschnitt 6 — Ergebnisse
# ==========================================================================
section_divider("06", "Ergebnisse",
                "Datenqualität vs. Algorithmus – die 4×4-Cross-Validation-Matrix")

# Heatmap-Folie
res = content_slide("Cross-Validation: F1 und ROC-AUC", "Ergebnisse")
res.shapes.add_picture(f1_png, Inches(0.55), Inches(1.85), height=Inches(4.55))
res.shapes.add_picture(auc_png, Inches(6.7), Inches(1.85), height=Inches(4.55))
_, rcf = textbox(res, Inches(0.6), Inches(6.6), Inches(12.0), Inches(0.5))
add_para(rcf, "Orange umrandet: Diagonale (trainiert = evaluiert). Zeilen = Trainingsquelle, Spalten = Evaluationsquelle.",
         12, MIDGRAY, italic=True)

# Kernbefunde
bullets_slide(
    "Kernbefunde der Matrix", "Interpretation",
    [
        "Datenqualität dominiert: routine (Aktiv-Rate ~26 %) erreicht auf sich selbst F1 0,83 und ROC-AUC 0,96.",
        "Synthetisch ≠ real: das emobpy-Modell generalisiert schwach (F1 nahe 0 auf den meisten Quellen).",
        "ved ist degeneriert: F1 = 0 selbst auf sich selbst – der Classifier sagt durchgängig „geparkt\" (Imbalance-Falle).",
        "Bester Transfer: real_world_ev → routine erreicht F1 0,66 – reale Trainingsdaten übertragen besser als synthetische.",
        ("Antwort auf FF1a-i: Die Datenquelle erklärt den Großteil der Streuung – mehr als die Wahl des Algorithmus.", 0),
    ])

# ==========================================================================
#  Abschnitt 7 — Diskussion
# ==========================================================================
section_divider("07", "Diskussion & Limitationen")

two_col_slide(
    "Einordnung & Grenzen", "Diskussion",
    "Implikationen",
    ["Quellenspezifische Modellierung ist gerechtfertigt (Heterogenität bestätigt).",
     "Explizite Unsicherheit (P10/P90) ist für eine sichere V2G-Steuerung essenziell.",
     "Datenqualitäts-Kriterien (Aktiv-Rate, Regelmäßigkeit) wirken direkt auf die Prognosegüte."],
    "Limitationen",
    ["Zeitliche Auflösung fix bei 1 Stunde (Datenpipeline-Diskretisierung).",
     "Imbalance verzerrt schwache Quellen (ved kollabiert).",
     "Ablation noch offen: LSTM / TFT als nächste Stufen ausstehend.",
     "Robustheit gegen unvollständige/fehlerhafte Daten noch in Arbeit."])

# ==========================================================================
#  Abschnitt 8 — Fazit
# ==========================================================================
section_divider("08", "Fazit & Ausblick")

# Beantwortung Forschungsfragen
fz = content_slide("Beantwortung der Forschungsfragen", "Fazit")
qa = [
    ("FF1", "Individuelles Verhalten ist über 7 Tage prognostizierbar – aber stark quellenabhängig. Regelmäßige reale Profile (routine) sind gut vorhersagbar, niedrigaktive Quellen schwierig."),
    ("FF1a/i", "Die Datenquelle erklärt mehr Varianz als der Algorithmus. RF und LightGBM liefern ähnliche Güte; die Datenqualität ist der dominante Hebel."),
    ("FF1b", "Zeit- und Lag-/Rolling-Merkmale (hour, weekday, Lags 1/24/168) tragen die Vorhersage."),
    ("FF2", "Qualität = ausreichende Aktiv-Rate + Regelmäßigkeit. Stündliche Diskretisierung statt künstlicher Interpolation hält die Daten ehrlich."),
]
yy = 1.85
for tag, txt in qa:
    rect(fz, Inches(0.62), Inches(yy), Inches(1.5), Inches(1.05), BLUE)
    _, tcf = textbox(fz, Inches(0.62), Inches(yy), Inches(1.5), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tcf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); set_run(r, tag, 14, WHITE, bold=True)
    _, bcf = textbox(fz, Inches(2.3), Inches(yy), Inches(10.2), Inches(1.05),
                     anchor=MSO_ANCHOR.MIDDLE)
    add_para(bcf, txt, 13.5, DARK, space_after=0, line_spacing=1.0)
    yy += 1.2

# Ausblick
bullets_slide(
    "Ausblick", "Nächste Schritte",
    [
        "Ablationsstudie vervollständigen: LogReg → RF → LightGBM → LSTM → TFT.",
        "Robustheit härten: Umgang mit unvollständigen, fehlerhaften und unregelmäßigen Daten.",
        "Imbalance-Strategien für niedrigaktive Quellen (Resampling, fokussierte Verluste).",
        "Anbindung an reale V2G-Steuerung und Validierung der Quantil-Kalibrierung im Betrieb.",
    ])

# Danke / Backup-Hinweis
end = prs.slides.add_slide(BLANK)
rect(end, 0, 0, EMU_W, EMU_H, NAVY)
rect(end, 0, 0, Inches(0.25), EMU_H, ORANGE)
rect(end, 0, Inches(4.3), EMU_W, Pt(2.5), BLUE)
_, ef = textbox(end, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6))
add_para(ef, "Vielen Dank für Ihre Aufmerksamkeit.", 32, WHITE, bold=True,
         font=FONT_LIGHT, space_after=8)
add_para(ef, "Ich freue mich auf Ihre Fragen.", 18, RGBColor(0xBF, 0xD8, 0xEE))
_, ef2 = textbox(end, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.0))
add_para(ef2, "Achim Baumgärtner   ·   Prognose individuellen Mobilitätsverhaltens   ·   Hochschule Esslingen",
         13, RGBColor(0xCE, 0xD9, 0xE3))

prs.save(OUT)
print("OK ->", OUT)
print("Folien:", len(prs.slides._sldIdLst))
